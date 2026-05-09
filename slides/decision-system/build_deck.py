"""Builds slides/decision-system/decision_system.pptx.

Slide 1 (executive): five native shapes laid out left-to-right with a memory
loop above, brand-coloured per the plan.
Slide 2 (technical): full-height embed of the rendered mermaid technical PNG.

Re-run:
    source .venv/bin/activate
    python slides/decision-system/build_deck.py
"""

from __future__ import annotations

import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


HERE = pathlib.Path(__file__).resolve().parent
OUT = HERE / "decision_system.pptx"
TECH_PNG = HERE / "technical.png"


# ---------- Brand palette (per plan) ----------
GREEN_PRIMARY = RGBColor(0x15, 0x80, 0x3D)        # NetIQ green deep
GREEN_BORDER = RGBColor(0x0B, 0x5E, 0x2A)
BRAIN_BORDER = RGBColor(0x15, 0x80, 0x3D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE_CAMARA_FILL = RGBColor(0xDB, 0xEA, 0xFE)
BLUE_CAMARA_BORDER = RGBColor(0x1D, 0x4E, 0xD8)
BLUE_CAMARA_TEXT = RGBColor(0x0A, 0x1F, 0x4D)
NEUTRAL_FILL = RGBColor(0xF1, 0xF5, 0xF9)
NEUTRAL_BORDER = RGBColor(0x64, 0x74, 0x8B)
NEUTRAL_TEXT = RGBColor(0x0F, 0x17, 0x2A)
ARROW_LINE = RGBColor(0x47, 0x55, 0x69)
TITLE_TEXT = RGBColor(0x0F, 0x17, 0x2A)
SUBTITLE_TEXT = RGBColor(0x47, 0x55, 0x69)


# ---------- Slide geometry (16:9, 13.333 x 7.5 in) ----------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_title(slide, text: str, subtitle: str | None = None) -> None:
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.13), Inches(0.7))
    tf = title.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = "Inter"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = TITLE_TEXT

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), Inches(12.13), Inches(0.4))
        stf = sub.text_frame
        stf.margin_left = stf.margin_right = stf.margin_top = stf.margin_bottom = 0
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.LEFT
        srun = sp.add_run()
        srun.text = subtitle
        srun.font.name = "Inter"
        srun.font.size = Pt(16)
        srun.font.italic = True
        srun.font.color.rgb = SUBTITLE_TEXT


def add_box(
    slide,
    *,
    left: Inches,
    top: Inches,
    width: Inches,
    height: Inches,
    headline: str,
    sub: str | None,
    fill: RGBColor,
    border: RGBColor,
    text_color: RGBColor,
    headline_size: int = 16,
    sub_size: int = 12,
    headline_bold: bool = True,
    shape: int = MSO_SHAPE.ROUNDED_RECTANGLE,
):
    box = slide.shapes.add_shape(shape, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = border
    box.line.width = Pt(1.5)
    box.shadow.inherit = False

    tf = box.text_frame
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = headline
    r1.font.name = "Inter"
    r1.font.size = Pt(headline_size)
    r1.font.bold = headline_bold
    r1.font.color.rgb = text_color

    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = sub
        r2.font.name = "Inter"
        r2.font.size = Pt(sub_size)
        r2.font.color.rgb = text_color

    return box


def connect(slide, src, dst, *, dashed: bool = False, label: str | None = None):
    """Draws a straight connector from `src.right_center` to `dst.left_center`
    (when dst is right of src) or vertically when stacked. Falls back to a
    straight line between centres if positions are arbitrary.
    """
    sx = src.left + src.width // 2
    sy = src.top + src.height // 2
    dx = dst.left + dst.width // 2
    dy = dst.top + dst.height // 2

    horizontal = abs(dy - sy) < abs(dx - sx)
    if horizontal:
        if dx > sx:
            sx = src.left + src.width
            dx = dst.left
        else:
            sx = src.left
            dx = dst.left + dst.width
        sy = src.top + src.height // 2
        dy = dst.top + dst.height // 2
    else:
        if dy > sy:
            sy = src.top + src.height
            dy = dst.top
        else:
            sy = src.top
            dy = dst.top + dst.height
        sx = src.left + src.width // 2
        dx = dst.left + dst.width // 2

    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, sx, sy, dx, dy)
    line.line.color.rgb = ARROW_LINE
    line.line.width = Pt(1.75)
    if dashed:
        # python-pptx exposes dash style via XML
        from pptx.oxml.ns import qn
        from lxml import etree
        ln = line.line._get_or_add_ln()
        prst_dash = etree.SubElement(ln, qn("a:prstDash"))
        prst_dash.set("val", "dash")
    # Add an end arrowhead
    from lxml import etree
    from pptx.oxml.ns import qn
    ln = line.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")

    if label:
        lbl_w = Inches(0.9)
        lbl_h = Inches(0.3)
        lbl_left = Emu((sx + dx) // 2 - lbl_w // 2)
        lbl_top = Emu((sy + dy) // 2 - lbl_h // 2)
        lbl = slide.shapes.add_textbox(lbl_left, lbl_top, lbl_w, lbl_h)
        ltf = lbl.text_frame
        ltf.margin_left = ltf.margin_right = ltf.margin_top = ltf.margin_bottom = 0
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        lrun = lp.add_run()
        lrun.text = label
        lrun.font.name = "Inter"
        lrun.font.size = Pt(10)
        lrun.font.italic = True
        lrun.font.color.rgb = ARROW_LINE


# ---------- Slide 1: Executive view ----------
def slide_executive(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    add_title(
        slide,
        "NetIQ decision making",
        "One decision, network-grade trust. Every call learns from every other sector.",
    )

    # Five-box left-to-right flow.
    # Vertical center y for the main row: ~4.4in
    row_y = Inches(3.95)
    row_h = Inches(1.4)

    # Memory box sits above the centre, between Brain and CAMARA
    mem_w, mem_h = Inches(2.4), Inches(1.1)

    # Compute layout
    margin_x = Inches(0.55)
    box_w = Inches(2.25)
    gap = Inches(0.45)

    request = add_box(
        slide,
        left=margin_x,
        top=row_y,
        width=box_w,
        height=row_h,
        headline="Request",
        sub="intent · phone · context",
        fill=GREEN_PRIMARY,
        border=GREEN_BORDER,
        text_color=WHITE,
    )

    brain_left = margin_x + box_w + gap
    brain = add_box(
        slide,
        left=brain_left,
        top=row_y,
        width=box_w,
        height=row_h,
        headline="Decision brain",
        sub="LLM agent or policy rules",
        fill=WHITE,
        border=BRAIN_BORDER,
        text_color=NEUTRAL_TEXT,
    )

    camara_left = brain_left + box_w + gap
    camara_top = row_y
    camara_h = row_h
    camara = add_box(
        slide,
        left=camara_left,
        top=camara_top,
        width=box_w,
        height=camara_h,
        headline="Nokia CAMARA APIs",
        sub="17 network signals",
        fill=BLUE_CAMARA_FILL,
        border=BLUE_CAMARA_BORDER,
        text_color=BLUE_CAMARA_TEXT,
    )

    decision_left = camara_left + box_w + gap
    decision_w = Inches(2.7)
    decision = add_box(
        slide,
        left=decision_left,
        top=row_y,
        width=decision_w,
        height=row_h,
        headline="Decision",
        sub="ALLOW · VERIFY · BLOCK\n+ confidence + reason",
        fill=GREEN_PRIMARY,
        border=GREEN_BORDER,
        text_color=WHITE,
    )

    # Memory box: above the brain/camara midpoint, tucked under title
    mem_left = brain_left + (box_w + gap + box_w - mem_w) // 2
    mem_top = Inches(1.85)
    memory = add_box(
        slide,
        left=mem_left,
        top=mem_top,
        width=mem_w,
        height=mem_h,
        headline="Cross-sector memory",
        sub="risk_profiles",
        fill=NEUTRAL_FILL,
        border=NEUTRAL_BORDER,
        text_color=NEUTRAL_TEXT,
    )

    # Connectors
    connect(slide, request, brain)
    connect(slide, brain, camara)
    connect(slide, camara, brain)  # back-edge: brain calls then receives
    connect(slide, brain, decision)
    connect(slide, memory, brain)

    # Learn back-edge: from decision (top centre) up and curve to memory (right side).
    # We approximate with a polyline using two connectors plus a label.
    from lxml import etree
    from pptx.oxml.ns import qn

    # Vertical leg up from decision top
    leg1_x = decision.left + decision.width // 2
    leg1_top = mem_top + mem_h // 2
    leg1_bot = decision.top
    leg1 = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, leg1_x, leg1_bot, leg1_x, leg1_top
    )
    leg1.line.color.rgb = ARROW_LINE
    leg1.line.width = Pt(1.5)
    ln1 = leg1.line._get_or_add_ln()
    dash1 = etree.SubElement(ln1, qn("a:prstDash"))
    dash1.set("val", "dash")

    # Horizontal leg from above-decision back to memory right edge
    mem_right_x = memory.left + memory.width
    leg2 = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, leg1_x, leg1_top, mem_right_x, leg1_top
    )
    leg2.line.color.rgb = ARROW_LINE
    leg2.line.width = Pt(1.5)
    ln2 = leg2.line._get_or_add_ln()
    dash2 = etree.SubElement(ln2, qn("a:prstDash"))
    dash2.set("val", "dash")
    tail = etree.SubElement(ln2, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")

    # "learn" label sitting on the horizontal leg
    lbl_w = Inches(0.9)
    lbl_h = Inches(0.32)
    lbl_left = Emu((mem_right_x + leg1_x) // 2 - lbl_w // 2)
    lbl_top = Emu(leg1_top - lbl_h // 2)
    lbl = slide.shapes.add_textbox(lbl_left, lbl_top, lbl_w, lbl_h)
    lbl.fill.solid()
    lbl.fill.fore_color.rgb = WHITE
    lbl.line.fill.background()
    p = lbl.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "learn"
    run.font.name = "Inter"
    run.font.size = Pt(11)
    run.font.italic = True
    run.font.color.rgb = ARROW_LINE

    # Footer caption
    foot = slide.shapes.add_textbox(Inches(0.6), Inches(6.65), Inches(12.13), Inches(0.35))
    ftf = foot.text_frame
    ftf.margin_left = ftf.margin_right = ftf.margin_top = ftf.margin_bottom = 0
    fp = ftf.paragraphs[0]
    fp.alignment = PP_ALIGN.LEFT
    fr = fp.add_run()
    fr.text = (
        "Mode: agent or policy. Trace, audit log, and memory update happen on every call."
    )
    fr.font.name = "Inter"
    fr.font.size = Pt(11)
    fr.font.italic = True
    fr.font.color.rgb = SUBTITLE_TEXT

    # Speaker notes
    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "A request arrives — could be from a fintech, a clinic, a co-op.\n"
        "We pull what NetIQ already knows about that line from cross-sector memory.\n"
        "The brain — either an LLM agent or a policy rule, your choice — picks the right "
        "Nokia CAMARA signals to call.\n"
        "We come back with ALLOW, VERIFY, or BLOCK in seconds, with confidence and reason.\n"
        "And every decision teaches the memory. That's the moat — the more sectors plug in, "
        "the smarter NetIQ gets for all of them."
    )


# ---------- Slide 2: Technical view (native shapes) ----------
def _add_section_label(slide, *, left, top, width, text, color):
    box = slide.shapes.add_textbox(left, top, width, Inches(0.28))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = "Inter"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = color


def slide_technical(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    add_title(
        slide,
        "How a NetIQ decision is made",
        "Dual-mode brain · agentic CAMARA orchestration · cross-sector memory loop.",
    )

    # Slide canvas usable area: x in [0.4, 12.93], y in [1.55, 7.1]
    # Three horizontal bands:
    #   Band A (top, y=1.7..3.05):   Request -> mode? -> two branches -> Decision
    #   Band B (middle, y=3.25..5.4): Brain box (4 nodes), Policy box, CAMARA box
    #   Band C (bottom, y=5.6..6.95): Memory + Audit, with feedback loop

    # ----- Band A: input row -----
    request = add_box(
        slide, left=Inches(0.5), top=Inches(1.75),
        width=Inches(2.1), height=Inches(0.95),
        headline="Request", sub="intent · phone · context · mode",
        fill=GREEN_PRIMARY, border=GREEN_BORDER, text_color=WHITE,
        headline_size=14, sub_size=10,
    )

    router = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND, Inches(2.95), Inches(1.7),
        Inches(1.3), Inches(1.05),
    )
    router.fill.solid()
    router.fill.fore_color.rgb = NEUTRAL_FILL
    router.line.color.rgb = NEUTRAL_BORDER
    router.line.width = Pt(1.5)
    router.shadow.inherit = False
    rtf = router.text_frame
    rtf.margin_left = rtf.margin_right = rtf.margin_top = rtf.margin_bottom = 0
    rtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    rp = rtf.paragraphs[0]
    rp.alignment = PP_ALIGN.CENTER
    rrun = rp.add_run()
    rrun.text = "mode?"
    rrun.font.name = "Inter"
    rrun.font.size = Pt(13)
    rrun.font.bold = True
    rrun.font.color.rgb = NEUTRAL_TEXT

    decision = add_box(
        slide, left=Inches(10.6), top=Inches(1.75),
        width=Inches(2.4), height=Inches(0.95),
        headline="Decision",
        sub="ALLOW · VERIFY · BLOCK · confidence · reason · trace",
        fill=GREEN_PRIMARY, border=GREEN_BORDER, text_color=WHITE,
        headline_size=14, sub_size=9.5,
    )

    # ----- Band B: dual-mode brains -----
    brain_left = Inches(4.6)
    brain_top = Inches(3.0)
    brain_w = Inches(4.5)
    brain_h = Inches(2.55)

    # Subgraph border (Agent brain wrapper)
    agent_wrap = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, brain_left, brain_top, brain_w, brain_h,
    )
    agent_wrap.fill.solid()
    agent_wrap.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    agent_wrap.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    agent_wrap.line.width = Pt(1)
    agent_wrap.shadow.inherit = False

    _add_section_label(
        slide,
        left=brain_left + Inches(0.12),
        top=brain_top + Inches(0.08),
        width=Inches(2.0),
        text="Agent brain",
        color=NEUTRAL_BORDER,
    )

    # 4 sub-nodes in a 2x2 grid inside the brain wrapper
    inner_pad = Inches(0.18)
    label_band = Inches(0.42)
    inner_left = brain_left + inner_pad
    inner_top = brain_top + label_band
    cell_w = (brain_w - 2 * inner_pad - Inches(0.18)) // 2
    cell_h = (brain_h - label_band - 2 * inner_pad - Inches(0.18)) // 2

    llm = add_box(
        slide, left=inner_left, top=inner_top,
        width=cell_w, height=cell_h,
        headline="LLM orchestrator",
        sub="gpt-4o-mini · tool-calling",
        fill=WHITE, border=BRAIN_BORDER, text_color=NEUTRAL_TEXT,
        headline_size=11, sub_size=9,
    )
    risk = add_box(
        slide, left=inner_left + cell_w + Inches(0.18), top=inner_top,
        width=cell_w, height=cell_h,
        headline="RiskAgent",
        sub="identity signals",
        fill=WHITE, border=BRAIN_BORDER, text_color=NEUTRAL_TEXT,
        headline_size=11, sub_size=9,
    )
    net = add_box(
        slide, left=inner_left, top=inner_top + cell_h + Inches(0.18),
        width=cell_w, height=cell_h,
        headline="NetworkAgent",
        sub="connectivity signals",
        fill=WHITE, border=BRAIN_BORDER, text_color=NEUTRAL_TEXT,
        headline_size=11, sub_size=9,
    )
    decide = add_box(
        slide, left=inner_left + cell_w + Inches(0.18),
        top=inner_top + cell_h + Inches(0.18),
        width=cell_w, height=cell_h,
        headline="DecisionAgent",
        sub="fusion + verdict",
        fill=WHITE, border=BRAIN_BORDER, text_color=NEUTRAL_TEXT,
        headline_size=11, sub_size=9,
    )

    # Policy engine sits below the brain wrapper, same column-band
    policy = add_box(
        slide, left=brain_left, top=brain_top + brain_h + Inches(0.18),
        width=brain_w, height=Inches(0.85),
        headline="Policy engine",
        sub="JSON rules · deterministic · audit-friendly",
        fill=RGBColor(0xFE, 0xF3, 0xC7),
        border=RGBColor(0xB4, 0x53, 0x09),
        text_color=RGBColor(0x3A, 0x22, 0x07),
        headline_size=13, sub_size=10,
    )

    # CAMARA cluster on the right side of brain
    camara_left = brain_left + brain_w + Inches(0.4)
    camara_top = brain_top
    camara_w = Inches(3.7)
    camara_h = brain_h
    camara_wrap = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, camara_left, camara_top, camara_w, camara_h,
    )
    camara_wrap.fill.solid()
    camara_wrap.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    camara_wrap.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    camara_wrap.line.width = Pt(1)
    camara_wrap.shadow.inherit = False

    _add_section_label(
        slide,
        left=camara_left + Inches(0.12),
        top=camara_top + Inches(0.08),
        width=Inches(3.5),
        text="Nokia Network as Code · CAMARA",
        color=BLUE_CAMARA_BORDER,
    )

    cam_pad = Inches(0.18)
    cam_band = Inches(0.42)
    cam_inner_left = camara_left + cam_pad
    cam_inner_top = camara_top + cam_band
    cam_cell_w = camara_w - 2 * cam_pad
    cam_cell_h = (camara_h - cam_band - 2 * cam_pad - Inches(0.36)) // 3

    cam_ident = add_box(
        slide, left=cam_inner_left, top=cam_inner_top,
        width=cam_cell_w, height=cam_cell_h,
        headline="Identity",
        sub="SIM swap · device swap · number verification · recycling",
        fill=BLUE_CAMARA_FILL, border=BLUE_CAMARA_BORDER, text_color=BLUE_CAMARA_TEXT,
        headline_size=11, sub_size=9,
    )
    cam_loc = add_box(
        slide, left=cam_inner_left, top=cam_inner_top + cam_cell_h + Inches(0.18),
        width=cam_cell_w, height=cam_cell_h,
        headline="Location",
        sub="verification · geofencing · roaming",
        fill=BLUE_CAMARA_FILL, border=BLUE_CAMARA_BORDER, text_color=BLUE_CAMARA_TEXT,
        headline_size=11, sub_size=9,
    )
    cam_qual = add_box(
        slide, left=cam_inner_left, top=cam_inner_top + 2 * (cam_cell_h + Inches(0.18)),
        width=cam_cell_w, height=cam_cell_h,
        headline="Quality",
        sub="QoS · reachability · congestion · consent",
        fill=BLUE_CAMARA_FILL, border=BLUE_CAMARA_BORDER, text_color=BLUE_CAMARA_TEXT,
        headline_size=11, sub_size=9,
    )

    # ----- Band C: memory + audit (left side of brain)
    memory = add_box(
        slide, left=Inches(0.5), top=Inches(3.0),
        width=Inches(3.7), height=Inches(1.05),
        headline="Cross-sector memory",
        sub="risk_profiles · learned from every prior decision",
        fill=NEUTRAL_FILL, border=NEUTRAL_BORDER, text_color=NEUTRAL_TEXT,
        headline_size=12, sub_size=10,
        shape=MSO_SHAPE.CAN,
    )
    audit = add_box(
        slide, left=Inches(0.5), top=Inches(4.5),
        width=Inches(3.7), height=Inches(1.05),
        headline="Audit log",
        sub="analyze_events · regulator-ready trace",
        fill=NEUTRAL_FILL, border=NEUTRAL_BORDER, text_color=NEUTRAL_TEXT,
        headline_size=12, sub_size=10,
        shape=MSO_SHAPE.CAN,
    )

    # ----- Connectors -----
    connect(slide, request, router)
    connect(slide, router, llm, label="agent")
    connect(slide, router, policy, label="policy")

    # Brain internal flow
    connect(slide, llm, risk)
    connect(slide, llm, net)
    connect(slide, risk, decide)
    connect(slide, net, decide)

    # Brain + policy reach into CAMARA
    connect(slide, risk, cam_ident)
    connect(slide, net, cam_qual)
    connect(slide, policy, cam_loc)

    # Brain + policy reach Decision
    connect(slide, decide, decision)
    connect(slide, policy, decision)

    # Memory feeds in to brain + policy
    connect(slide, memory, llm)
    connect(slide, memory, policy)

    # Decision -> Audit -> dashed back to Memory
    connect(slide, decision, audit)
    connect(slide, audit, memory, dashed=True, label="memory update")

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "Technical view of the same pipeline. The mode router splits agent vs policy. "
        "On the agent path the LLM orchestrator picks among RiskAgent (identity signals) "
        "and NetworkAgent (connectivity signals); both call into the Nokia CAMARA tool layer "
        "(identity, location, quality). DecisionAgent fuses the signals into an "
        "ALLOW / VERIFY / BLOCK verdict with confidence, reason, and a full trace. "
        "Every decision is persisted into analyze_events for audit and feeds back into "
        "cross-sector memory for the next call."
    )


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_executive(prs)
    slide_technical(prs)

    prs.save(OUT)
    print(f"Wrote {OUT} ({OUT.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
